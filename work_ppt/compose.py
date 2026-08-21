from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches

from work_ppt.diagrams import FLOW_KINDS, apply_diagram
from work_ppt.gate import GateError
from work_ppt.mermaid import add_mermaid, export_drawio
from work_ppt.onboard import SKIP_PH, onboard
from work_ppt.qa import QaError


def _delete_slides(prs: Presentation) -> None:
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _layout_by_name(prs: Presentation, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"Layout not found: {name}")


def pick_layout(profile: dict, hint: str, slot_count: int) -> str:
    names = [l["name"] for l in profile["layouts"]]
    if hint in names:
        return hint
    aliases = {
        "cover": ["title-cover", "Title Slide"],
        "section": ["title-centered", "Section Header"],
        "title-body": ["content-centered-a", "content-centered-b", "Title and Content"],
        "two-col": ["column-2-centered", "Two Content", "Comparison"],
        "three-col": ["column-3-centered-a", "column-3", "column-3-centered-b", "column-3-centered-c"],
        "four-col": ["column-4-centered"],
        "blank": ["Blank", "master-base"],
    }
    for cand in aliases.get(hint, []):
        if cand in names:
            return cand
    # Weak master: only title + content
    if "Title and Content" in names:
        return "Title and Content"
    if "Title Slide" in names and hint == "cover":
        return "Title Slide"
    return names[0]


def _insert_picture(slide, image: Path) -> None:
    image = Path(image)
    if not image.is_file():
        raise GateError(f"picture not found: {image}")
    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            shape.insert_picture(str(image))
            return
    slide.shapes.add_picture(str(image), Inches(0.6), Inches(2.4), width=Inches(5.5))


def _fill_placeholders(slide, slots: list[str]) -> None:
    fillable = []
    for shape in slide.placeholders:
        if shape.placeholder_format.type in SKIP_PH:
            continue
        if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            continue
        fillable.append(shape)
    fillable.sort(key=lambda s: (0 if "TITLE" in str(s.placeholder_format.type) else 1, s.placeholder_format.idx))
    for shape, text in zip(fillable, slots):
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        tf.clear()
        lines = text.split("\n")
        p0 = tf.paragraphs[0]
        p0.text = lines[0]
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0 if not line.startswith("- ") else 0


def compose(template: Path, plan: dict, dest: Path) -> Path:
    template = Path(template)
    dest = Path(dest)
    dest.parent.mkdir(parents=True, exist_ok=True)
    profile = onboard(template)
    prs = Presentation(str(template))
    _delete_slides(prs)
    for slide_spec in plan["slides"]:
        slots = slide_spec.get("slots") or []
        hint = slide_spec.get("layout_hint") or "title-body"
        layout_name = pick_layout(profile, hint, len(slots))
        layout = _layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(layout)
        _fill_placeholders(slide, slots)
        slide_spec["resolved_layout"] = layout_name
        kind = slide_spec.get("diagram")
        gen = slide_spec.get("generate_image")
        if gen and kind in FLOW_KINDS:
            raise GateError(
                "generated images cannot replace process/architecture/sequence diagrams"
            )
        if kind:
            apply_diagram(slide, kind, slide_spec.get("diagram_labels") or [])
        picture = slide_spec.get("picture")
        if not picture and isinstance(gen, str) and gen not in {"true", "atmosphere"}:
            picture = gen
        if picture:
            _insert_picture(slide, Path(picture))
        drawio = slide_spec.get("drawio")
        if drawio:
            src = Path(drawio)
            if src.suffix.lower() == ".svg":
                _insert_picture(slide, src)
            else:
                svg = dest.parent / f"{dest.stem}-slide-drawio.svg"
                _insert_picture(slide, export_drawio(src, svg))
    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))
    for index, slide_spec in enumerate(plan["slides"], start=1):
        mermaid = slide_spec.get("mermaid")
        if not mermaid:
            continue
        try:
            add_mermaid(dest, index, mermaid)
        except QaError:
            raise GateError(
                "officecli mermaid failed; install officecli or drop the mermaid field"
            )
    return dest


def load_plan(path: Path) -> dict:
    return json.loads(Path(path).read_text(encoding="utf-8"))
