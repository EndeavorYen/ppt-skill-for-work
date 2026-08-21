from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x00, 0x44, 0x79)
INK = RGBColor(0x11, 0x11, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _box(slide, x, y, w, h, text, fill=NAVY, font_color=WHITE, size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = font_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def _arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = INK
    conn.line.width = Pt(1.5)
    return conn


def add_process_row(slide, labels: list[str], top=Inches(2.1), connect=True) -> None:
    n = len(labels)
    if n == 0:
        return
    margin = Inches(0.5)
    gap = Inches(0.25)
    width = Inches(12.3)
    box_w = int((width - gap * (n - 1)) / n)
    box_h = Inches(1.05)
    x = margin
    boxes = []
    for label in labels:
        boxes.append(_box(slide, x, top, box_w, box_h, label, size=11))
        x += box_w + gap
    if not connect:
        return
    for a, b in zip(boxes, boxes[1:]):
        _arrow(
            slide,
            a.left + a.width,
            a.top + a.height // 2,
            b.left,
            b.top + b.height // 2,
        )


def add_sibling_row(slide, labels: list[str], top=Inches(3.2)) -> None:
    add_process_row(slide, labels, top=top, connect=False)


def add_architecture_stack(slide, layers: list[str], left=Inches(1.2), top=Inches(1.8)) -> None:
    w = Inches(10.5)
    h = Inches(0.85)
    gap = Inches(0.18)
    fills = [RGBColor(0x00, 0x44, 0x79), RGBColor(0x00, 0x05, 0x42), RGBColor(0x5E, 0x5E, 0x5E), NAVY]
    y = top
    for i, layer in enumerate(layers):
        _box(slide, left, y, w, h, layer, fill=fills[i % len(fills)], size=14)
        y += h + gap


def add_csa_hca_fork(slide) -> None:
    """Near-window feeds parallel CSA and HCA — not a three-layer stack."""
    src = _box(
        slide,
        Inches(3.4),
        Inches(2.15),
        Inches(6.5),
        Inches(0.95),
        "近窗 128 token 未壓縮 KV（兩條路徑都保留）",
        size=13,
    )
    left = _box(
        slide,
        Inches(0.7),
        Inches(4.15),
        Inches(5.5),
        Inches(1.2),
        "CSA：m=4 + 稀疏 top-k\n保細節",
        fill=RGBColor(0x00, 0x05, 0x42),
        size=13,
    )
    right = _box(
        slide,
        Inches(7.1),
        Inches(4.15),
        Inches(5.5),
        Inches(1.2),
        "HCA：m'=128 再 dense\n保覆蓋",
        fill=RGBColor(0x5E, 0x5E, 0x5E),
        size=13,
    )
    _arrow(slide, src.left + src.width // 2, src.top + src.height, left.left + left.width // 2, left.top)
    _arrow(slide, src.left + src.width // 2, src.top + src.height, right.left + right.width // 2, right.top)


FLOW_KINDS = {
    "process",
    "sequence",
    "architecture",
    "sibling",
    "sibling-costs",
    "decoder-callouts",
    "csa-hca-fork",
}


def apply_diagram(slide, kind: str, labels: list[str] | None = None) -> None:
    labels = list(labels or [])
    if kind in {"sibling-costs"}:
        add_sibling_row(
            slide, labels or ["KV 記憶體", "Attention FLOPs", "Residual 表達力"]
        )
    elif kind in {"sibling"}:
        add_sibling_row(slide, labels)
    elif kind in {"process", "sequence"}:
        add_process_row(slide, labels)
    elif kind in {"architecture"}:
        add_architecture_stack(slide, labels)
    elif kind == "decoder-callouts":
        add_decoder_callouts(slide)
    elif kind == "csa-hca-fork":
        add_csa_hca_fork(slide)
    else:
        raise ValueError(f"unknown diagram {kind}")


def add_decoder_callouts(slide) -> None:
    core = _box(
        slide,
        Inches(3.6),
        Inches(2.35),
        Inches(6.0),
        Inches(1.7),
        "Decoder block\nAttention + FFN",
        size=16,
    )
    callouts = [
        (Inches(0.45), Inches(1.85), "KV-share\n跨層重用 K/V"),
        (Inches(9.55), Inches(1.85), "PLE\nlookup 補容量"),
        (Inches(0.45), Inches(4.55), "CCA\n潛空間算 attention"),
        (Inches(9.55), Inches(4.55), "mHC\nn=4 residual"),
    ]
    boxes = []
    for x, y, text in callouts:
        boxes.append(_box(slide, x, y, Inches(3.0), Inches(1.15), text, fill=RGBColor(0x00, 0x05, 0x42), size=12))
    cx = core.left + core.width // 2
    cy = core.top + core.height // 2
    for box in boxes:
        bx = box.left + box.width // 2
        by = box.top + box.height // 2
        _arrow(slide, bx, by, cx, cy)
