from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

MASTER_NOISE = (
    "click to edit",
    "click to add",
    "edit master",
    "drag picture",
    "click icon",
    "second level",
    "full name",
    "job position",
)


def _is_noise(text: str) -> bool:
    low = text.strip().lower()
    if not low or low in {"‹#›", "#"}:
        return True
    return any(n in low for n in MASTER_NOISE)


def extract(deck: Path) -> dict:
    deck = Path(deck)
    prs = Presentation(str(deck))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        blocks = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = "\n".join(
                p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
            )
            if _is_noise(text):
                continue
            role = "body"
            if shape.is_placeholder:
                t = str(shape.placeholder_format.type)
                if "TITLE" in t:
                    role = "title"
                elif shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                    continue
            blocks.append({"role": role, "text": text, "shape": shape.name})
        slides.append(
            {
                "index": i,
                "layout": slide.slide_layout.name,
                "blocks": blocks,
            }
        )
    return {
        "source": str(deck).replace("\\", "/"),
        "slide_count": len(slides),
        "slides": slides,
    }


def save_extract(data: dict, dest: Path) -> Path:
    dest = Path(dest)
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
    return dest
