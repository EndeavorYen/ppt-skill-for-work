from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from work_ppt.compose import _delete_slides, _fill_placeholders, _layout_by_name, compose
from work_ppt.extract import extract
from work_ppt.gate import (
    assert_numbers_sourced,
    load_json,
    number_tokens,
    source_corpus,
    story_to_plan,
)
from work_ppt.onboard import onboard

REPO = Path(__file__).resolve().parents[1]
TEMPLATE = REPO / "docs/fixtures/templates/dense-consulting-inner-chapter.pptx"
GOLD = REPO / "eval/gold"
FIXTURES = REPO / "docs/fixtures/templates"

TEMPLATES = {
    "inner": TEMPLATE,
    "weak": FIXTURES / "impoverished-title-content.pptx",
    "light": FIXTURES / "light-corporate-office-default.pptx",
    "dark": FIXTURES / "dark-tech-navy.pptx",
}

CASE_TABLE = {
    "ab": {"family": "inner", "kind": "ab"},
    "sourced": {"family": "inner", "kind": "sourced"},
    "weak-ab": {"family": "weak", "kind": "ab"},
    "light-ab": {"family": "light", "kind": "ab"},
    "dark-ab": {"family": "dark", "kind": "ab"},
    "light-sourced": {"family": "light", "kind": "sourced"},
    "dark-sourced": {"family": "dark", "kind": "sourced"},
    "weak-sourced": {"family": "weak", "kind": "sourced"},
}
CASES = tuple(CASE_TABLE)
CASE_ORDER = (
    "ab",
    "sourced",
    "weak-ab",
    "light-ab",
    "dark-ab",
    "light-sourced",
    "dark-sourced",
    "weak-sourced",
)

CATALOG_STORY = {
    "title": "LLM Architecture Update",
    "decision": "None — catalog closer",
    "slides": [
        {
            "action_title": "LLM Architecture Update",
            "layout_hint": "cover",
            "slots": [
                "LLM Architecture Update",
                "Gemma 4 / Laguna / ZAYA1 / DeepSeek V4",
                "Internal sharing",
            ],
        },
        {
            "action_title": "Agenda",
            "layout_hint": "title-body",
            "slots": [
                "Agenda",
                "Introduction\nGemma 4\nLaguna XS.2\nZAYA1-8B\nDeepSeek V4\nConclusion",
            ],
        },
        {
            "action_title": "Background",
            "layout_hint": "title-body",
            "slots": [
                "Background",
                "LLMs are getting better\nContext is longer\nThere are many new models in 2026\nArchitecture is important",
            ],
        },
        {
            "action_title": "Gemma 4 Overview",
            "layout_hint": "two-col",
            "slots": [
                "Gemma 4 Overview",
                "E2B and E4B for devices\n26B MoE\n31B dense\nUses GQA",
                "KV sharing across layers\nPLE embeddings\nSaves memory\nCode on GitHub",
            ],
        },
        {
            "action_title": "Laguna XS.2",
            "layout_hint": "title-body",
            "slots": [
                "Laguna XS.2",
                "Poolside coding model\n40 layers\nSliding window + global attention\nDifferent query heads per layer\nSimilar to OpenELM",
            ],
        },
        {
            "action_title": "ZAYA1-8B",
            "layout_hint": "title-body",
            "slots": [
                "ZAYA1-8B",
                "Zyphra model on AMD GPUs\nCompressed Convolutional Attention\nRelated to MLA\nMoE with one expert\nPaper on arXiv",
            ],
        },
        {
            "action_title": "DeepSeek V4",
            "layout_hint": "title-body",
            "slots": [
                "DeepSeek V4",
                "Biggest release this year\nmHC residual streams\nCSA and HCA compression\nBetter than V3.2\nVery sparse MoE",
            ],
        },
        {
            "action_title": "Key Numbers",
            "layout_hint": "title-body",
            "slots": [
                "Key Numbers",
                "Some memory savings\nLong context is cheaper\nLots of complexity\nNeed to keep learning",
            ],
        },
        {
            "action_title": "Thank You",
            "layout_hint": "section",
            "slots": ["Thank You", "Questions?"],
        },
    ],
}


def extract_blob(data: dict) -> str:
    return "\n".join(
        block["text"]
        for slide in data.get("slides") or []
        for block in slide.get("blocks") or []
    )


def story_blob(story: dict) -> str:
    parts = []
    for slide in story.get("slides") or []:
        parts.append(slide.get("action_title") or "")
        parts.extend(slide.get("slots") or [])
        parts.extend(slide.get("diagram_labels") or [])
    return "\n".join(parts)


def case_dir(case: str) -> Path:
    spec = CASE_TABLE.get(case)
    if not spec:
        raise ValueError(f"unknown gold case {case}; expected {CASES}")
    if spec["family"] == "inner":
        return GOLD / spec["kind"]
    return GOLD / spec["family"] / spec["kind"]


def family_original_path(family: str) -> Path:
    if family == "inner":
        return GOLD / "original.pptx"
    return GOLD / family / "original.pptx"


def build_original(dest: Path) -> Path:
    """Typical AI-first-draft on Inner Chapter layouts."""
    dest = Path(dest)
    dest.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE))
    _delete_slides(prs)

    def add(layout_name, slots):
        slide = prs.slides.add_slide(_layout_by_name(prs, layout_name))
        _fill_placeholders(slide, slots)

    add("title-cover", ["LLM Architecture Update", "Gemma 4 / Laguna / ZAYA1 / DeepSeek V4", "Internal sharing"])
    add(
        "content-centered-a",
        [
            "Agenda",
            "Introduction\nGemma 4\nLaguna XS.2\nZAYA1-8B\nDeepSeek V4\nConclusion",
        ],
    )
    add(
        "content-centered-a",
        [
            "Background",
            "LLMs are getting better\nContext is longer\nThere are many new models in 2026\nArchitecture is important",
        ],
    )
    add(
        "column-2-centered",
        [
            "Gemma 4 Overview",
            "E2B and E4B for devices\n26B MoE\n31B dense\nUses GQA",
            "KV sharing across layers\nPLE embeddings\nSaves memory\nCode on GitHub",
        ],
    )
    add(
        "content-centered-a",
        [
            "Laguna XS.2",
            "Poolside coding model\n40 layers\nSliding window + global attention\nDifferent query heads per layer\nSimilar to OpenELM",
        ],
    )
    add(
        "content-centered-a",
        [
            "ZAYA1-8B",
            "Zyphra model on AMD GPUs\nCompressed Convolutional Attention\nRelated to MLA\nMoE with one expert\nPaper on arXiv",
        ],
    )
    add(
        "content-centered-a",
        [
            "DeepSeek V4",
            "Biggest release this year\nmHC residual streams\nCSA and HCA compression\nBetter than V3.2\nVery sparse MoE",
        ],
    )
    add(
        "content-centered-a",
        [
            "Key Numbers",
            "Some memory savings\nLong context is cheaper\nLots of complexity\nNeed to keep learning",
        ],
    )
    add("title-centered", ["Thank You", "Questions?"])
    prs.save(str(dest))
    return dest


def build_family_original(family: str, dest: Path | None = None) -> Path:
    dest = Path(dest) if dest else family_original_path(family)
    dest.parent.mkdir(parents=True, exist_ok=True)
    if family == "inner":
        return build_original(dest)
    template = TEMPLATES[family]
    plan = story_to_plan(CATALOG_STORY, onboard(template))
    return compose(template, plan, dest)


def assert_ab_locked(original_extract: dict, story: dict) -> None:
    allowed = set(number_tokens(extract_blob(original_extract)))
    found = set(number_tokens(story_blob(story)))
    extra = sorted(found - allowed, key=lambda token: (len(token), token))
    if extra:
        raise ValueError(f"ab story introduces numbers absent from original: {extra}")


def _takeaway_titles(extract_data: dict) -> bool:
    titles = [
        block["text"]
        for slide in extract_data["slides"]
        for block in slide["blocks"]
        if block["role"] == "title"
    ]
    if not titles:
        titles = [
            slide["blocks"][0]["text"]
            for slide in extract_data["slides"]
            if slide["blocks"]
        ]
    banned = {"Agenda", "Background", "Thank You", "Key Numbers"}
    return all(title not in banned and len(title) >= 12 for title in titles)


def gold_review(
    original: Path,
    optimized: Path,
    dest: Path,
    family: str = "inner",
) -> dict:
    src = extract(original)
    out = extract(optimized)
    profile = onboard(original)
    allowed_layouts = {layout["name"] for layout in profile["layouts"]}
    out_layouts = {slide["layout"] for slide in out["slides"]}
    same_template = out_layouts <= allowed_layouts
    titles = " ".join(
        block["text"]
        for slide in out["slides"]
        for block in slide["blocks"]
        if block["role"] == "title"
    )
    no_thank_you = "Thank You" not in titles
    has_native_shapes = any(
        "Rounded Rectangle" in (block.get("shape") or "")
        for slide in out["slides"]
        for block in slide["blocks"]
    )
    slide_count_gte = out["slide_count"] >= src["slide_count"]
    takeaways = _takeaway_titles(out)
    structural = {
        "no_thank_you": no_thank_you,
        "slide_count_gte": slide_count_gte,
        "has_native_shapes": has_native_shapes,
        "takeaway_titles": takeaways,
        "no_invented_master": same_template,
    }
    if family == "weak":
        passed = same_template and no_thank_you and takeaways
    else:
        passed = same_template and no_thank_you and slide_count_gte and has_native_shapes
    payload = {
        "same_template": same_template,
        "family": family,
        "structural": structural,
        "pass": passed,
    }
    dest = Path(dest)
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")
    return payload


def build_case(case: str, original: Path | None = None, dest: Path | None = None) -> Path:
    spec = CASE_TABLE[case]
    family = spec["family"]
    kind = spec["kind"]
    folder = case_dir(case)
    folder.mkdir(parents=True, exist_ok=True)
    if family == "inner" and original is not None:
        orig = Path(original)
    else:
        orig = build_family_original(family, original if original else None)
    story = load_json(GOLD / kind / "story.json")
    brief = dict(load_json(GOLD / kind / "brief.json"))
    brief["prior_deck"] = str(orig).replace("\\", "/")
    brief["template"] = str(TEMPLATES[family]).replace("\\", "/")
    extracted = extract(orig)
    corpus = source_corpus(brief, json.dumps(extracted, ensure_ascii=False))
    assert_numbers_sourced(story, corpus)
    if kind == "ab":
        fact_src = extract(GOLD / "original.pptx") if (GOLD / "original.pptx").exists() else extracted
        assert_ab_locked(fact_src, story)
    plan = story_to_plan(story, onboard(orig))
    dest = Path(dest) if dest else folder / "optimized.pptx"
    compose(orig, plan, dest)
    (folder / "brief.json").write_text(
        json.dumps(brief, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    (dest.with_suffix(".plan.json")).write_text(
        json.dumps(plan, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    return dest


def build_optimized(dest: Path, template: Path | None = None) -> Path:
    original = Path(template) if template else GOLD / "original.pptx"
    if not original.exists():
        original = build_original(Path(dest).parent / "original.pptx")
    return build_case("sourced", original, dest)
