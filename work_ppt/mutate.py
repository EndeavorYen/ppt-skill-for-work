from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation

from work_ppt.compose import _fill_placeholders, pick_layout
from work_ppt.extract import extract
from work_ppt.gate import GateError
from work_ppt.onboard import onboard

MUTATE_LIMIT = 0.30


def changed_fraction(extract_data: dict, plan: dict) -> float:
    n = extract_data["slide_count"]
    if n == 0:
        return 1.0
    changed = 0
    current = extract_data["slides"]
    planned = plan.get("slides") or []
    for i, spec in enumerate(planned):
        if i >= n:
            changed += 1
            continue
        if spec.get("resolved_layout") != current[i]["layout"]:
            changed += 1
    if len(planned) > n:
        changed += len(planned) - n
    return changed / n


def mutate(deck: Path, plan: dict, dest: Path) -> Path:
    deck = Path(deck)
    dest = Path(dest)
    dest.parent.mkdir(parents=True, exist_ok=True)
    extract_data = extract(deck)
    profile = onboard(deck)
    for spec in plan.get("slides") or []:
        if not spec.get("resolved_layout"):
            hint = spec.get("layout_hint") or "title-body"
            slots = spec.get("slots") or []
            spec["resolved_layout"] = pick_layout(profile, hint, len(slots))
    fraction = changed_fraction(extract_data, plan)
    if fraction > MUTATE_LIMIT:
        raise GateError(
            f"mutate would change {fraction:.0%} of slides; use optimize"
        )
    if len(plan.get("slides") or []) > extract_data["slide_count"]:
        raise GateError("mutate cannot add slides; use optimize")
    shutil.copy2(deck, dest)
    prs = Presentation(str(dest))
    for i, spec in enumerate(plan.get("slides") or []):
        _fill_placeholders(prs.slides[i], spec.get("slots") or [])
    prs.save(str(dest))
    return dest
