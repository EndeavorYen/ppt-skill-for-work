from __future__ import annotations

import json
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

SKIP_PH = {
    PP_PLACEHOLDER.SLIDE_NUMBER,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.HEADER,
}


def _officecli_theme(path: Path) -> dict:
    try:
        raw = subprocess.check_output(
            ["officecli", "get", str(path), "/", "--json"],
            text=True,
            stderr=subprocess.STDOUT,
        )
        data = json.loads(raw)
        fmt = (data.get("data") or {}).get("results") or [{}]
        return fmt[0].get("format") or {}
    except (OSError, subprocess.CalledProcessError, json.JSONDecodeError, IndexError):
        return {}


def onboard(template: Path) -> dict:
    template = Path(template)
    prs = Presentation(str(template))
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for shape in layout.placeholders:
            ph = shape.placeholder_format
            placeholders.append(
                {
                    "idx": int(ph.idx),
                    "type": str(ph.type),
                    "name": shape.name,
                    "fillable": ph.type not in SKIP_PH,
                }
            )
        fillable = [p for p in placeholders if p["fillable"]]
        layouts.append(
            {
                "index": i,
                "name": layout.name,
                "placeholder_count": len(placeholders),
                "body_slots": sum(1 for p in fillable if "BODY" in p["type"] or "OBJECT" in p["type"]),
                "has_title": any("TITLE" in p["type"] for p in fillable),
                "placeholders": placeholders,
            }
        )
    theme = _officecli_theme(template)
    profile = {
        "template": str(template).replace("\\", "/"),
        "slide_width_emu": int(prs.slide_width),
        "slide_height_emu": int(prs.slide_height),
        "layout_count": len(layouts),
        "theme": {
            "name": theme.get("theme.name"),
            "dk1": theme.get("theme.color.dk1"),
            "lt1": theme.get("theme.color.lt1"),
            "accent1": theme.get("theme.color.accent1"),
            "accent2": theme.get("theme.color.accent2"),
            "major_font": theme.get("theme.font.major.latin"),
            "minor_font": theme.get("theme.font.minor.latin"),
        },
        "layouts": layouts,
    }
    if not layouts:
        raise ValueError(f"No slide layouts in {template}")
    return profile


def save_profile(profile: dict, dest: Path) -> Path:
    dest = Path(dest)
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_text(json.dumps(profile, indent=2, ensure_ascii=False), encoding="utf-8")
    return dest
